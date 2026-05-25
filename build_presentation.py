"""
build_presentation.py
=====================
Builds FINM3422 A3 — Multi-Asset Risk & Derivatives Platform
15-slide PowerPoint presentation.

Output: docs/A3_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
import sys
import io
import tempfile

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from scipy.stats import norm as scipy_norm
from scipy.interpolate import CubicSpline

sys.path.insert(0, '/home/user/FINM3422_A3_T7_G1_DP')

# ─────────────────────────────────────────────────────────────────────────────
# CHART GENERATION
# ─────────────────────────────────────────────────────────────────────────────

def generate_charts(chart_dir):
    """Generate all 7 chart PNGs into chart_dir."""
    os.makedirs(chart_dir, exist_ok=True)

    # ── Load data ──────────────────────────────────────────────────────────
    f17 = pd.read_csv('/home/user/FINM3422_A3_T7_G1_DP/data/F17_DATA_CLEAN.csv')
    eq_ret = pd.read_csv('/home/user/FINM3422_A3_T7_G1_DP/data/equity_returns.csv')

    # Last row of F17 — zero coupon rates
    last_row = f17.iloc[-1]
    maturities = np.array([float(c) for c in f17.columns[1:]])
    rates_last = np.array([float(last_row[c]) for c in f17.columns[1:]])

    # ── CHART 1: Yield Curve ───────────────────────────────────────────────
    cs = CubicSpline(maturities, rates_last)
    T_fine = np.linspace(0, 10, 400)
    r_linear = np.interp(T_fine, maturities, rates_last)
    r_cubic = cs(T_fine)
    D_linear = np.exp(-r_linear * T_fine)
    D_obs = np.exp(-rates_last * maturities)

    fig, axes = plt.subplots(1, 2, figsize=(16, 5))
    ax = axes[0]
    ax.plot(T_fine, r_linear * 100, 'b-', linewidth=1.8, label='Linear')
    ax.plot(T_fine, r_cubic * 100, 'r--', linewidth=1.8, label='Cubic Spline')
    ax.scatter(maturities, rates_last * 100, color='black', s=20, zorder=5, label='Observed Rates')
    ax.set_title('Zero-Coupon Yield Curve (RBA F17, 30 Apr 2026)', fontsize=12, fontweight='bold')
    ax.set_xlabel('Maturity (years)')
    ax.set_ylabel('Zero Rate (%)')
    ax.legend(fontsize=9)
    ax.grid(True, alpha=0.4)
    ax.set_xlim(0, 10)

    ax2 = axes[1]
    ax2.plot(T_fine, D_linear, 'g-', linewidth=2, label='Linear D(T)')
    ax2.scatter(maturities, D_obs, color='black', s=20, zorder=5)
    ax2.set_title('Discount Factors D(T)', fontsize=12, fontweight='bold')
    ax2.set_xlabel('Maturity (years)')
    ax2.set_ylabel('Discount Factor D(T)')
    ax2.legend(fontsize=9)
    ax2.grid(True, alpha=0.4)
    ax2.set_xlim(0, 10)

    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart1_yield_curve.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart1_yield_curve.png saved")

    # ── CHART 2: Binomial Convergence ─────────────────────────────────────
    S0 = 164.29
    K  = 164.29 * 0.95
    T  = 0.5
    sigma = 0.2337
    q  = 0.0299
    r  = 0.0446   # from yield curve at T=0.5

    # BS put price
    d1 = (np.log(S0 / K) + (r - q + 0.5 * sigma**2) * T) / (sigma * np.sqrt(T))
    d2 = d1 - sigma * np.sqrt(T)
    bs_price = K * np.exp(-r * T) * scipy_norm.cdf(-d2) - S0 * np.exp(-q * T) * scipy_norm.cdf(-d1)

    steps = [10, 20, 50, 100, 200, 500, 1000, 2000]
    bin_prices = []
    for N in steps:
        dt   = T / N
        u    = np.exp(sigma * np.sqrt(dt))
        d    = 1.0 / u
        disc = np.exp(-r * dt)
        p    = (np.exp((r - q) * dt) - d) / (u - d)
        # terminal payoffs
        j    = np.arange(N + 1)
        ST   = S0 * (u ** (N - j)) * (d ** j)
        V    = np.maximum(K - ST, 0.0)
        # backward induction
        for _ in range(N):
            V = disc * (p * V[:-1] + (1 - p) * V[1:])
        bin_prices.append(float(V[0]))

    fig, ax = plt.subplots(figsize=(10, 5))
    ax.plot(steps, bin_prices, 'b-o', linewidth=1.8, markersize=5, label='Binomial Price')
    ax.axhline(bs_price, color='red', linestyle='--', linewidth=1.8, label=f'BS Price (${bs_price:.2f})')
    ax.annotate(f'BS = ${bs_price:.2f}', xy=(steps[-1], bs_price),
                xytext=(steps[-1] * 0.6, bs_price + 0.05),
                fontsize=10, color='red',
                arrowprops=dict(arrowstyle='->', color='red'))
    ax.set_title('Binomial Tree Convergence to Black-Scholes — CBA 6m OTM Put', fontsize=12, fontweight='bold')
    ax.set_xlabel('Tree Steps (N)')
    ax.set_ylabel('Option Price ($)')
    ax.legend(fontsize=10)
    ax.grid(True, alpha=0.4)
    ax.set_xscale('log')
    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart2_option_convergence.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart2_option_convergence.png saved")

    # ── CHART 3: Dollar-Delta Bar Chart ────────────────────────────────────
    tickers  = ['CBA', 'BHP', 'CSL']
    dd_vals  = [66443, 58338, 102757]
    colors3  = ['#2E86AB', '#A23B72', '#E07A5F']
    labels3  = ['Protective Put\n(Hedging)', 'Covered Call\n(Income)', 'Long Straddle\n(Speculation)']

    fig, ax = plt.subplots(figsize=(9, 5))
    bars = ax.bar(tickers, dd_vals, color=colors3, width=0.5, edgecolor='white', linewidth=1.5)
    for bar, val in zip(bars, dd_vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 800,
                f'${val:,.0f}', ha='center', va='bottom', fontsize=11, fontweight='bold')
    ax.set_xticks(range(len(tickers)))
    ax.set_xticklabels([f'{t}\n{l}' for t, l in zip(tickers, labels3)], fontsize=10)
    ax.set_title('Dollar-Delta Exposure by Strategy', fontsize=13, fontweight='bold')
    ax.set_ylabel('Dollar-Delta ($)')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${x:,.0f}'))
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_ylim(0, max(dd_vals) * 1.18)
    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart3_dollar_delta.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart3_dollar_delta.png saved")

    # ── CHART 4: P&L Distribution ──────────────────────────────────────────
    weights  = np.array([0.292, 0.256, 0.452])
    exposure = 227538.0
    port_ret = (eq_ret['CBA'] * weights[0] +
                eq_ret['BHP'] * weights[1] +
                eq_ret['CSL'] * weights[2]).dropna()
    pnl = port_ret.values * exposure

    # VaR / ES
    alpha = 0.05
    hist_var = -np.percentile(pnl, alpha * 100)
    mu_p = pnl.mean()
    sig_p = pnl.std()
    param_var = -(mu_p + scipy_norm.ppf(alpha) * sig_p)
    es_vals = pnl[pnl <= -hist_var]
    hist_es = -es_vals.mean() if len(es_vals) else hist_var

    fig, ax = plt.subplots(figsize=(11, 5))
    ax.hist(pnl, bins=60, color='steelblue', density=True, alpha=0.75, label='Daily P&L')
    x_range = np.linspace(pnl.min(), pnl.max(), 400)
    ax.plot(x_range, scipy_norm.pdf(x_range, mu_p, sig_p), 'coral', linewidth=2, label='Normal fit')
    ax.axvline(-hist_var, color='red', linestyle='--', linewidth=1.8,
               label=f'Hist VaR 95% (${hist_var:,.0f})')
    ax.axvline(-param_var, color='orange', linestyle=':', linewidth=2,
               label=f'Param VaR 95% (${param_var:,.0f})')
    ax.axvline(-hist_es, color='darkred', linestyle='-.', linewidth=2,
               label=f'ES 95% (${hist_es:,.0f})')
    ax.set_title('Portfolio Daily P&L Distribution\nHistorical vs Parametric VaR at 95% confidence',
                 fontsize=11, fontweight='bold')
    ax.set_xlabel('Daily P&L ($)')
    ax.set_ylabel('Density')
    ax.legend(fontsize=9)
    ax.grid(True, alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart4_pnl_distribution.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart4_pnl_distribution.png saved")

    # ── CHART 5: Correlation Stress ────────────────────────────────────────
    regimes_labels = ['Base\n(ρ≈0.07)', 'Stressed\n(ρ≈0.85)', 'Crisis\n(ρ≈0.99)']
    var_vals5 = [12775, 17927, 18379]
    es_vals5  = [15599, 21787, 22675]
    x5 = np.arange(len(regimes_labels))
    width5 = 0.35

    fig, ax = plt.subplots(figsize=(10, 5))
    b1 = ax.bar(x5 - width5/2, var_vals5, width5, label='VaR (95%, 10-day)',
                color=['green', 'orange', 'red'], edgecolor='white')
    b2 = ax.bar(x5 + width5/2, es_vals5,  width5, label='ES (95%, 10-day)',
                color=['green', 'orange', 'red'], edgecolor='white', hatch='///', alpha=0.85)
    for bar in list(b1) + list(b2):
        ax.text(bar.get_x() + bar.get_width() / 2,
                bar.get_height() + 150,
                f'${bar.get_height():,.0f}', ha='center', va='bottom', fontsize=8)
    ax.set_xticks(x5)
    ax.set_xticklabels(regimes_labels, fontsize=11)
    ax.set_title('VaR & ES Across Correlation Regimes (95%, 10-day)', fontsize=12, fontweight='bold')
    ax.set_ylabel('Value ($)')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${x:,.0f}'))
    ax.legend(fontsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_ylim(0, max(es_vals5) * 1.2)
    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart5_correlation_stress.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart5_correlation_stress.png saved")

    # ── CHART 6: Scenario Analysis ─────────────────────────────────────────
    scenarios6  = ['Base case', 'Equity −10%', 'Equity −20%',
                   'Equity +10%', 'Equity +20%', 'Vol +30%', 'Vol +50%',
                   'Crash (−15%+50%vol)']
    pnl6 = [0, -21782, -41180, 23402, 47565, 2150, 3587, -22000]
    colors6 = ['#4CAF50' if v >= 0 else '#FF6B6B' for v in pnl6]

    fig, ax = plt.subplots(figsize=(11, 5))
    bars6 = ax.barh(scenarios6, pnl6, color=colors6, edgecolor='white', linewidth=1)
    for bar, val in zip(bars6, pnl6):
        offset = 300 if val >= 0 else -300
        ha = 'left' if val >= 0 else 'right'
        ax.text(val + offset, bar.get_y() + bar.get_height() / 2,
                f'${val:+,.0f}', ha=ha, va='center', fontsize=9, fontweight='bold')
    ax.axvline(0, color='black', linewidth=0.8)
    ax.set_title('Portfolio P&L Under Stress Scenarios', fontsize=13, fontweight='bold')
    ax.set_xlabel('Portfolio P&L ($)')
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${x:+,.0f}'))
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, axis='x', alpha=0.3)
    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart6_scenarios.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart6_scenarios.png saved")

    # ── CHART 7: Drawdown ──────────────────────────────────────────────────
    w7 = np.array([0.292, 0.256, 0.452])
    port_ret7 = (eq_ret['CBA'] * w7[0] +
                 eq_ret['BHP'] * w7[1] +
                 eq_ret['CSL'] * w7[2]).dropna().values
    exposure7 = 227538.0
    cum_wealth = exposure7 * np.exp(np.cumsum(port_ret7))
    running_peak = np.maximum.accumulate(cum_wealth)
    drawdown = (cum_wealth - running_peak) / running_peak * 100  # percentage
    idx = np.arange(len(cum_wealth))

    mdd = drawdown.min()

    fig, axes7 = plt.subplots(2, 1, figsize=(11, 7), sharex=True,
                               gridspec_kw={'height_ratios': [2, 1]})

    ax_top = axes7[0]
    ax_top.plot(idx, cum_wealth, 'b-', linewidth=1.5, label='Portfolio Value')
    ax_top.plot(idx, running_peak, 'purple', linestyle='--', linewidth=1.2, label='Running Peak')
    ax_top.fill_between(idx, cum_wealth, running_peak, where=(cum_wealth < running_peak),
                        color='red', alpha=0.15, label='Drawdown region')
    ax_top.annotate(f'MDD = {mdd:.2f}%',
                    xy=(idx[np.argmin(cum_wealth)], cum_wealth.min()),
                    xytext=(idx[np.argmin(cum_wealth)] + 20, cum_wealth.min() * 1.05),
                    fontsize=9, color='darkred',
                    arrowprops=dict(arrowstyle='->', color='darkred'))
    ax_top.set_title('Full Portfolio Wealth Path & Drawdown (Dollar-Delta Weighted)',
                     fontsize=11, fontweight='bold')
    ax_top.set_ylabel('Portfolio Value ($)')
    ax_top.legend(fontsize=9)
    ax_top.grid(True, alpha=0.3)
    ax_top.spines['top'].set_visible(False)
    ax_top.spines['right'].set_visible(False)
    ax_top.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'${x:,.0f}'))

    ax_bot = axes7[1]
    ax_bot.plot(idx, drawdown, 'r-', linewidth=1.2)
    ax_bot.fill_between(idx, drawdown, 0, where=(drawdown < 0), color='red', alpha=0.3)
    ax_bot.set_title('Drawdown Series', fontsize=10, fontweight='bold')
    ax_bot.set_ylabel('Drawdown (%)')
    ax_bot.set_xlabel('Trading Day Index')
    ax_bot.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{x:.1f}%'))
    ax_bot.grid(True, alpha=0.3)
    ax_bot.spines['top'].set_visible(False)
    ax_bot.spines['right'].set_visible(False)

    fig.tight_layout()
    fig.savefig(os.path.join(chart_dir, 'chart7_drawdown.png'), dpi=150, bbox_inches='tight')
    plt.close(fig)
    print("  chart7_drawdown.png saved")

    print("All charts generated.")
    return chart_dir


# ── Colour palette ────────────────────────────────────────────────────────────
C_BG_TITLE   = RGBColor(0xEB, 0xED, 0xF9)   # #EBEDF9 – title slide BG
C_BG_CONTENT = RGBColor(0xEE, 0xF2, 0xFF)   # #EEF2FF – content slide BG
C_BG_LAST    = RGBColor(0x1E, 0x27, 0x61)   # #1E2761 – last slide BG
C_NAVY       = RGBColor(0x2D, 0x3A, 0x6B)   # #2D3A6B – header bars / dark text
C_ACCENT     = RGBColor(0x5B, 0x8D, 0xEF)   # #5B8DEF – accent blue
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xFF)   # table alt-row tint

# ── Slide dimensions: 10" × 5.625" ──────────────────────────────────────────
W = Inches(10)
H = Inches(5.625)

OUT_PATH = "/home/user/FINM3422_A3_T7_G1_DP/docs/A3_Presentation.pptx"

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    """Add a completely blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]   # 'Blank' layout
    return prs.slides.add_slide(blank_layout)


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()   # no border
    else:
        shape.line.color.rgb = fill_color
    return shape


def add_textbox(slide, left, top, width, height,
                text, font_size, bold=False, italic=False,
                color=C_BLACK, align=PP_ALIGN.LEFT,
                word_wrap=True, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_textbox(slide, left, top, width, height,
                          lines,            # list of (text, size, bold, italic, color, align)
                          word_wrap=True, font_name='Calibri'):
    """Each element of `lines` is a tuple (text, size, bold, italic, color, align).
       Each element creates a separate paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    for i, (text, size, bold, italic, color, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        # support newlines within a single entry
        sub_lines = text.split('\n')
        for j, sub in enumerate(sub_lines):
            if j == 0:
                run = p.add_run()
                run.text = sub
                run.font.name = font_name
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = color
            else:
                p2 = tf.add_paragraph()
                p2.alignment = align
                run2 = p2.add_run()
                run2.text = sub
                run2.font.name = font_name
                run2.font.size = Pt(size)
                run2.font.bold = bold
                run2.font.italic = italic
                run2.font.color.rgb = color
    return txBox


def add_header_bar(slide, title, subtitle='', title_size=22, sub_size=11):
    """Standard dark-navy header bar across the top of a content slide."""
    bar = add_rect(slide,
                   left=Inches(0), top=Inches(0),
                   width=Inches(10), height=Inches(1.0),
                   fill_color=C_NAVY)

    # Title inside bar
    tb_title = slide.shapes.add_textbox(Inches(0.3), Inches(0.05),
                                         Inches(9.4), Inches(0.55))
    tf = tb_title.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Calibri'
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = C_WHITE

    if subtitle:
        tb_sub = slide.shapes.add_textbox(Inches(0.3), Inches(0.60),
                                           Inches(9.4), Inches(0.35))
        tf2 = tb_sub.text_frame
        tf2.word_wrap = False
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = 'Calibri'
        r2.font.size = Pt(sub_size)
        r2.font.italic = True
        r2.font.color.rgb = RGBColor(0xCC, 0xD6, 0xFF)


def add_slide_number(slide, n):
    tb = slide.shapes.add_textbox(Inches(9.2), Inches(5.3), Inches(0.7), Inches(0.25))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(n)
    r.font.name = 'Calibri'
    r.font.size = Pt(9)
    r.font.color.rgb = RGBColor(0x99, 0xAA, 0xCC)


def style_cell(cell, bg_color=None, font_color=C_BLACK,
               font_size=9, bold=False, italic=False,
               align=PP_ALIGN.CENTER, font_name='Calibri'):
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    cell.text_frame.word_wrap = True
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name  = font_name
            run.font.size  = Pt(font_size)
            run.font.bold  = bold
            run.font.italic = italic
            run.font.color.rgb = font_color


def set_cell_text(cell, text, bg_color=None, font_color=C_BLACK,
                  font_size=9, bold=False, align=PP_ALIGN.CENTER,
                  font_name='Calibri'):
    cell.text = text
    style_cell(cell, bg_color=bg_color, font_color=font_color,
               font_size=font_size, bold=bold, align=align,
               font_name=font_name)


def add_content_box(slide, left, top, width, height,
                    header_text, body_lines,
                    header_bg=C_NAVY, body_bg=C_LIGHT_GREY,
                    header_font_size=12, body_font_size=9):
    """A rounded box with a coloured header band and text body."""
    # Header rectangle
    add_rect(slide, left, top, width, Inches(0.35), header_bg)
    # Body rectangle
    add_rect(slide, left, top + Inches(0.35), width,
             height - Inches(0.35), body_bg)
    # Header text
    tb_h = slide.shapes.add_textbox(left + Inches(0.05), top + Inches(0.03),
                                     width - Inches(0.1), Inches(0.3))
    tf = tb_h.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = header_text
    r.font.name  = 'Calibri'
    r.font.size  = Pt(header_font_size)
    r.font.bold  = True
    r.font.color.rgb = C_WHITE
    # Body text
    body_text = '\n'.join(body_lines)
    tb_b = slide.shapes.add_textbox(left + Inches(0.08),
                                     top + Inches(0.38),
                                     width - Inches(0.12),
                                     height - Inches(0.45))
    tf2 = tb_b.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    r2 = p2.add_run()
    r2.text = body_text
    r2.font.name  = 'Calibri'
    r2.font.size  = Pt(body_font_size)
    r2.font.color.rgb = C_NAVY


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title Slide
# ─────────────────────────────────────────────────────────────────────────────
def slide1(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_TITLE)

    # Right-side bars
    add_rect(sl, Inches(7.2), Inches(0), Inches(2.8), H, C_NAVY)
    add_rect(sl, Inches(8.8), Inches(0), Inches(1.2), H, C_ACCENT)

    # Sub-heading line
    add_textbox(sl, Inches(0.55), Inches(1.0), Inches(6.4), Inches(0.4),
                "FINM3422 — Tutorial 7, Group 1",
                14, bold=True, color=C_NAVY)

    # Main title (two lines)
    txBox = sl.shapes.add_textbox(Inches(0.55), Inches(1.55), Inches(6.4), Inches(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Multi-Asset Risk &\nDerivatives Platform"
    r.font.name  = 'Calibri'
    r.font.size  = Pt(36)
    r.font.bold  = True
    r.font.color.rgb = C_BG_LAST

    # Assignment line
    add_textbox(sl, Inches(0.55), Inches(3.05), Inches(6.4), Inches(0.45),
                "A3 — Risk & Derivatives Platform",
                12, italic=True, color=C_ACCENT)

    # Authors
    add_textbox(sl, Inches(0.55), Inches(3.65), Inches(6.4), Inches(0.975),
                "Indigo Hill – 47057465\nJordan Westerberg – 48917258\nElizabeth Kvyatkovska – 48426983\nMath Al-Abudi",
                11, color=C_NAVY)

    # Date
    add_textbox(sl, Inches(0.55), Inches(4.95), Inches(3.0), Inches(0.3),
                "May 2026", 11, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Platform Overview
# ─────────────────────────────────────────────────────────────────────────────
def slide2(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Platform Overview",
                   "How the three modules connect")
    add_slide_number(sl, 2)

    box_w = Inches(3.0)
    box_h = Inches(3.6)
    box_y = Inches(1.1)
    gap   = Inches(0.1)
    starts = [Inches(0.2), Inches(3.4), Inches(6.6)]
    headers = ["Yield Curve", "Derivatives Pricing", "Portfolio Risk"]
    bodies = [
        ["RBA F17 data",
         "41 maturities (0–10yr)",
         "Linear & cubic interpolation",
         "Continuous compounding",
         "Risk-free rate r(T)"],
        ["Black-Scholes (Merton)",
         "Binomial CRR tree",
         "Monte Carlo (100k paths)",
         "American early exercise",
         "Closed-form Greeks"],
        ["3 strategies (hedge/income/spec)",
         "Historical & Parametric VaR",
         "Multi-asset correlated MC VaR",
         "Scenario analysis",
         "Max Drawdown"],
    ]
    bg_colors = [C_NAVY, C_ACCENT, C_NAVY]

    for i in range(3):
        add_content_box(sl, starts[i], box_y, box_w, box_h,
                        headers[i], bodies[i],
                        header_bg=bg_colors[i],
                        body_bg=RGBColor(0xE8, 0xED, 0xFF) if i != 1 else RGBColor(0xD6, 0xE6, 0xFF),
                        header_font_size=13, body_font_size=10)

    # Arrows between boxes
    add_textbox(sl, Inches(3.2), Inches(2.6), Inches(0.2), Inches(0.3),
                "▶", 14, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(6.4), Inches(2.6), Inches(0.2), Inches(0.3),
                "▶", 14, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Bottom caption
    add_textbox(sl, Inches(0.2), Inches(4.85), Inches(9.6), Inches(0.3),
                "Platform workflow: yield curve → derivative pricing → portfolio risk analytics",
                9, italic=True, color=C_NAVY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Yield Curve Construction
# ─────────────────────────────────────────────────────────────────────────────
def slide3(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Yield Curve Construction",
                   "RBA F17 Zero-Coupon Rates — as of 30 April 2026")
    add_slide_number(sl, 3)

    # Narrow info panels on the left (4.5")
    left_w  = Inches(4.4)
    panel_y = Inches(1.1)
    panel_h = Inches(2.0)

    add_content_box(sl, Inches(0.2), panel_y, left_w, panel_h,
                    "Data Source",
                    ["Source: RBA F17 zero-coupon rates",
                     "Date: 30-04-2026  |  41 maturities (0–10yr, 0.25yr steps)",
                     "Period: 2017-01-03 to 2026-04-30"],
                    body_font_size=9)

    # Rate table below data source panel
    tbl_y = panel_y + panel_h + Inches(0.08)
    add_content_box(sl, Inches(0.2), tbl_y, left_w, Inches(2.0),
                    "Risk-Free Rates r (continuous, from zero-coupon strip)",
                    [],
                    body_font_size=9)
    tbl = sl.shapes.add_table(5, 3,
                               Inches(0.25), tbl_y + Inches(0.38),
                               left_w - Inches(0.1), Inches(1.55)).table
    rows_data = [
        ("Maturity", "r (%)", "D(T)"),
        ("3 months", "4.30",  "0.9893"),
        ("6 months", "4.46",  "0.9779"),
        ("1 year",   "4.63",  "0.9548"),
        ("2 years",  "4.69",  "0.9105"),
    ]
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            bg = C_NAVY if ri == 0 else (C_LIGHT_GREY if ri % 2 == 1 else C_WHITE)
            fc = C_WHITE if ri == 0 else C_NAVY
            set_cell_text(cell, val, bg_color=bg, font_color=fc,
                          font_size=9, bold=(ri == 0))

    # Chart on the right (5")
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart1_yield_curve.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(4.7), Inches(1.1),
                                  Inches(5.1), Inches(4.1))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Yield Curve: Interpolation Comparison
# ─────────────────────────────────────────────────────────────────────────────
def slide4(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Yield Curve Construction",
                   "Linear vs Cubic Spline — Interpolation Comparison")
    add_slide_number(sl, 4)

    # Explanatory text
    add_textbox(sl, Inches(0.25), Inches(1.08), Inches(9.5), Inches(0.35),
                "Comparison at 7 non-observed maturities. Maximum difference: 0.56bp — immaterial for 6-month option pricing.",
                10, color=C_NAVY)

    # Table
    headers = ["Maturity (yr)", "Linear Rate (%)", "Cubic Rate (%)", "Difference (bp)"]
    data = [
        ("0.6",  "4.408", "4.404", "+0.04"),
        ("1.3",  "4.641", "4.643", "−0.02"),
        ("2.8",  "4.712", "4.715", "−0.03"),
        ("4.2",  "4.762", "4.766", "−0.04"),
        ("6.3",  "4.803", "4.808", "−0.05"),
        ("8.5",  "4.831", "4.837", "−0.06"),
        ("9.1",  "4.843", "4.843", "+0.00"),
    ]

    tbl_left = Inches(0.25)
    tbl_top  = Inches(1.55)
    tbl_w    = Inches(9.5)
    tbl_h    = Inches(3.4)

    tbl = sl.shapes.add_table(len(data) + 1, 4,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_w = [Inches(2.0), Inches(2.5), Inches(2.5), Inches(2.5)]
    for ci, cw in enumerate(col_w):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=10, bold=True)

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
            set_cell_text(cell, val, bg_color=bg, font_color=C_NAVY,
                          font_size=10)

    # Bottom note
    add_textbox(sl, Inches(0.25), Inches(5.1), Inches(9.5), Inches(0.3),
                "Cubic spline selected for long-dated instruments; linear sufficient for 6m equity options.",
                9, italic=True, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Derivatives Pricing
# ─────────────────────────────────────────────────────────────────────────────
def slide5(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Derivatives Pricing",
                   "Black-Scholes · Binomial CRR · Monte Carlo · American Tree")
    add_slide_number(sl, 5)

    headers = ["Position", "Spot($)", "σ(%pa)", "q(%)",
               "BS Price($)", "Binomial($)", "MC 100k($)",
               "American($)", "Am. Prem.(%)"]
    data = [
        ("CBA 6m 5%OTM put",  "164.29", "23.37", "2.99", "6.39", "6.39", "6.39", "6.46", "1.15"),
        ("BHP 6m 5%OTM call", "60.35",  "24.95", "3.28", "3.09", "3.09", "3.09", "3.09", "−0.02"),
        ("CSL 6m ATM call",   "98.40",  "30.91", "4.27", "8.42", "8.42", "8.43", "8.45", "0.30"),
        ("CSL 6m ATM put",    "98.40",  "30.91", "4.27", "8.33", "8.33", "8.33", "8.37", "0.45"),
    ]

    tbl_left = Inches(0.15)
    tbl_top  = Inches(1.08)
    tbl_w    = Inches(9.7)
    tbl_h    = Inches(1.55)

    tbl = sl.shapes.add_table(len(data) + 1, 9,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(2.1), Inches(0.75), Inches(0.75), Inches(0.6),
                  Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=8, bold=True)

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=C_NAVY,
                          font_size=8, align=align)

    # Convergence chart (main body)
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart2_option_convergence.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(0.15), Inches(2.7),
                                  Inches(9.7), Inches(2.75))
    else:
        # Observations as fallback
        obs = [
            ("1.  Merton Extension:", "Continuous dividend yield q applied to all ASX stocks. Without this, calls are overpriced and puts underpriced."),
            ("2.  American Premium:", "Small (<1.2%) at 6-month maturity — European pricing retained. American pricer available for longer maturities."),
            ("3.  MC Validation:", "All four positions agree with BS to within $0.01 at 100k paths with antithetic variates."),
        ]
        y_start = Inches(2.95)
        for label, body in obs:
            tb_l = sl.shapes.add_textbox(Inches(0.2), y_start, Inches(2.0), Inches(0.52))
            tf = tb_l.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = label
            r.font.name = 'Calibri'
            r.font.size = Pt(9)
            r.font.bold = True
            r.font.color.rgb = C_ACCENT
            add_textbox(sl, Inches(2.2), y_start, Inches(7.6), Inches(0.52),
                        body, 9, color=C_NAVY)
            y_start += Inches(0.58)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Option Greeks
# ─────────────────────────────────────────────────────────────────────────────
def slide6(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Option Greeks",
                   "Closed-Form (Black-Scholes) per Portfolio Position")
    add_slide_number(sl, 6)

    headers = ["Position", "Δ Delta", "Γ Gamma", "V Vega($)",
               "Θ Theta($/day)", "ρ Rho($)"]
    data = [
        ("CBA 6m OTM put",  "−0.326", "0.0132", "4.81",  "−1.79", "−0.68"),
        ("BHP 6m OTM call", "0.431",       "0.0365", "1.69",  "−1.04", "0.40"),
        ("CSL 6m ATM call", "0.534",       "0.0180", "16.82", "−5.64", "0.70"),
        ("CSL 6m ATM put",  "−0.447", "0.0180", "16.65", "−4.49", "−0.58"),
    ]

    tbl_left = Inches(0.2)
    tbl_top  = Inches(1.1)
    tbl_w    = Inches(9.6)
    tbl_h    = Inches(1.65)

    tbl = sl.shapes.add_table(len(data) + 1, 6,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(2.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.9), Inches(1.5)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=10, bold=True)

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=C_NAVY,
                          font_size=10, align=align)

    # Portfolio aggregates
    agg_y = Inches(2.88)
    add_rect(sl, Inches(0.2), agg_y, Inches(9.6), Inches(1.0),
             RGBColor(0xD8, 0xE4, 0xFF))
    add_textbox(sl, Inches(0.3), agg_y + Inches(0.05),
                Inches(9.4), Inches(0.28),
                "Portfolio Net Aggregates",
                11, bold=True, color=C_NAVY)
    agg_text = ("Δ = +2,415 (equity-equivalent exposure)    "
                "Γ = −36.0 (net short gamma — short BHP call dominates)    "
                "V = +$238 per 1% vol (net long vega — CSL straddle drives this)\n"
                "Θ = −$16/day (modest time decay)    "
                "ρ = −$395 per 1% rate move")
    add_textbox(sl, Inches(0.3), agg_y + Inches(0.32),
                Inches(9.4), Inches(0.62),
                agg_text, 9.5, color=C_NAVY)

    # Key insight box
    ki_y = Inches(3.97)
    add_rect(sl, Inches(0.2), ki_y, Inches(9.6), Inches(0.75),
             RGBColor(0x2D, 0x3A, 0x6B))
    add_textbox(sl, Inches(0.3), ki_y + Inches(0.06),
                Inches(1.5), Inches(0.28),
                "Key Insight:", 10, bold=True, color=C_ACCENT)
    add_textbox(sl, Inches(0.3), ki_y + Inches(0.32),
                Inches(9.3), Inches(0.38),
                "FD Greeks sit on the Derivative base class using type(self); "
                "every subclass (Binomial, MC) automatically inherits a complete set of numerical Greeks.",
                9, color=C_WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Portfolio Construction
# ─────────────────────────────────────────────────────────────────────────────
def slide7(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Portfolio Construction",
                   "Three Strategies — Hedging · Income · Speculation")
    add_slide_number(sl, 7)

    headers = ["Position", "Qty", "Unit Val($)", "Pos. Val($)", "Unit Δ", "Pos. Δ"]
    data = [
        ("CBA equity (long, hedged)",        "600",    "164.30", "98,577",    "1.000",   "600.00"),
        ("CBA 6m 5% OTM put (protective)",   "600",    "6.39",   "3,834",     "−0.326", "−195.58"),
        ("BHP equity (long, covered)",       "1,700",  "60.35",  "102,587",   "1.000",   "1,700.00"),
        ("BHP 6m 5% OTM call (short)",       "−1,700", "3.09", "−5,251", "−0.431", "−733.26"),
        ("CSL equity (long, core)",          "1,000",  "98.40",  "98,400",    "1.000",   "1,000.00"),
        ("CSL 6m ATM call (straddle)",       "500",    "8.42",   "4,212",     "0.534",   "267.00"),
        ("CSL 6m ATM put (straddle)",        "500",    "8.33",   "4,165",     "−0.447", "−223.50"),
        ("TOTAL",                            "",       "",       "306,524",   "",        "2,415.43"),
    ]

    tbl_left = Inches(0.15)
    tbl_top  = Inches(1.08)
    tbl_w    = Inches(9.7)
    tbl_h    = Inches(2.5)

    tbl = sl.shapes.add_table(len(data) + 1, 6,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(3.0), Inches(0.9), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.0)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=9, bold=True)

    for ri, row in enumerate(data):
        is_total = ri == len(data) - 1
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if is_total:
                bg = C_NAVY
                fc = C_WHITE
                bd = True
            else:
                bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
                fc = C_NAVY
                bd = False
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=fc,
                          font_size=9, bold=bd, align=align)

    # Strategy summary
    strat_y = Inches(3.7)
    add_rect(sl, Inches(0.15), strat_y, Inches(9.7), Inches(1.6),
             RGBColor(0xE8, 0xED, 0xFF))

    strategies = [
        ("CBA:",
         "Long equity + protective put → hedging. Reduces directional exposure by ~$32k"),
        ("BHP:",
         "Long equity + short covered call → income. Caps upside, collects $5,251 premium"),
        ("CSL:",
         "Long equity + ATM straddle → vol speculation. Net delta ≈+0.09, pure vega bet"),
    ]
    for i, (lbl, body) in enumerate(strategies):
        row_y = strat_y + Inches(0.08) + Inches(i * 0.48)
        tb_l = sl.shapes.add_textbox(Inches(0.25), row_y, Inches(0.55), Inches(0.4))
        tf = tb_l.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = lbl
        r.font.name = 'Calibri'
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C_ACCENT
        add_textbox(sl, Inches(0.8), row_y, Inches(9.0), Inches(0.4),
                    body, 10, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Portfolio Exposures
# ─────────────────────────────────────────────────────────────────────────────
def slide8(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Portfolio Risk Exposures",
                   "Dollar-Delta Decomposition by Strategy")
    add_slide_number(sl, 8)

    # Exposure table (compact, top-left)
    headers = ["Ticker", "Strategy", "Dollar-Delta($)", "% of Total"]
    data = [
        ("CBA",   "Hedging (protective put)",    "66,443",  "29.2%"),
        ("BHP",   "Income (covered call)",        "58,338",  "25.6%"),
        ("CSL",   "Speculation (long straddle)", "102,757", "45.2%"),
        ("TOTAL", "",                             "227,538", "100.0%"),
    ]

    tbl_left = Inches(0.2)
    tbl_top  = Inches(1.1)
    tbl_w    = Inches(4.0)
    tbl_h    = Inches(1.45)

    tbl = sl.shapes.add_table(len(data) + 1, 4,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(0.6), Inches(1.8), Inches(1.0), Inches(0.6)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=9, bold=True)

    for ri, row in enumerate(data):
        is_total = ri == len(data) - 1
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if is_total:
                bg = C_NAVY; fc = C_WHITE; bd = True
            else:
                bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
                fc = C_NAVY; bd = False
            align = PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=fc,
                          font_size=9, bold=bd, align=align)

    # Dollar-delta chart (right side)
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart3_dollar_delta.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(4.35), Inches(1.05),
                                  Inches(5.5), Inches(4.35))

    # Key observations (below the table, left side)
    obs_y = Inches(2.65)
    add_rect(sl, Inches(0.2), obs_y, Inches(4.0), Inches(2.65),
             RGBColor(0xE0, 0xE8, 0xFF))
    add_textbox(sl, Inches(0.3), obs_y + Inches(0.05),
                Inches(3.8), Inches(0.28),
                "Key Observations", 10, bold=True, color=C_NAVY)
    obs_items = [
        "Total equity: $299,564",
        "Net option premium: −$6,961",
        "Portfolio vol: 18.22% (annualised)",
        "CBA put reduces $99k → $66k delta",
        "BHP short call caps upside",
        "CSL straddle = pure vol view",
    ]
    obs_text = '\n'.join(f"• {o}" for o in obs_items)
    add_textbox(sl, Inches(0.3), obs_y + Inches(0.36),
                Inches(3.8), Inches(2.2),
                obs_text, 9, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Value at Risk
# ─────────────────────────────────────────────────────────────────────────────
def slide9(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Value at Risk",
                   "Historical · Parametric · Multi-Asset MC — 95% confidence, 10-day horizon")
    add_slide_number(sl, 9)

    # P&L distribution chart on the left 65%
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart4_pnl_distribution.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(0.15), Inches(1.08),
                                  Inches(6.3), Inches(4.35))

    # VaR table on the right 35%
    headers = ["Method", "95% VaR ($)"]
    data = [
        ("Historical VaR",            "11,825"),
        ("Parametric VaR",            "13,583"),
        ("MC correlated",             "12,775"),
        ("ES (Hist)",                 "21,431"),
        ("ES (MC)",                   "15,599"),
    ]

    tbl_left = Inches(6.55)
    tbl_top  = Inches(1.1)
    tbl_w    = Inches(3.3)
    tbl_h    = Inches(2.1)

    tbl = sl.shapes.add_table(len(data) + 1, 2,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(1.95), Inches(1.35)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=9, bold=True)

    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
            align = PP_ALIGN.CENTER if ci == 1 else PP_ALIGN.LEFT
            set_cell_text(cell, val, bg_color=bg, font_color=C_NAVY,
                          font_size=9, align=align)

    # Key note (right side below table)
    note_y = Inches(3.3)
    add_rect(sl, Inches(6.55), note_y, Inches(3.3), Inches(2.1),
             RGBColor(0x2D, 0x3A, 0x6B))
    add_textbox(sl, Inches(6.65), note_y + Inches(0.06),
                Inches(3.1), Inches(0.28),
                "Interpretation:", 9, bold=True, color=C_ACCENT)
    add_textbox(sl, Inches(6.65), note_y + Inches(0.34),
                Inches(3.1), Inches(1.65),
                "Hist VaR < Param VaR — real returns are leptokurtic (fat tails). "
                "Empirical 5th pctl is more extreme than −1.645σ implies.\n\n"
                "Ann. Vol: 18.22%\nSkewness: Negative\nKurtosis: Leptokurtic",
                8.5, color=C_WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Leptokurtosis & VaR Interpretation
# ─────────────────────────────────────────────────────────────────────────────
def slide10(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Value at Risk — Leptokurtosis Effect",
                   "Why Historical VaR < Parametric VaR")
    add_slide_number(sl, 10)

    col_w = Inches(4.7)
    left_x = Inches(0.2)
    right_x = Inches(5.1)
    content_y = Inches(1.1)
    content_h = Inches(4.35)

    # LEFT column background
    add_rect(sl, left_x, content_y, col_w, content_h, RGBColor(0xE8, 0xEE, 0xFF))

    add_textbox(sl, left_x + Inches(0.1), content_y + Inches(0.08),
                col_w - Inches(0.2), Inches(0.35),
                "Return Distribution", 12, bold=True, color=C_NAVY)
    add_textbox(sl, left_x + Inches(0.1), content_y + Inches(0.42),
                col_w - Inches(0.2), Inches(0.5),
                "Daily P&L distribution has fatter tails than a normal distribution predicts.",
                10, color=C_NAVY)

    # Stats mini-table
    stat_top = content_y + Inches(0.95)
    stat_tbl = sl.shapes.add_table(5, 2,
                                    left_x + Inches(0.1),
                                    stat_top,
                                    col_w - Inches(0.2),
                                    Inches(1.5)).table
    stat_data = [
        ("Metric", "Value"),
        ("Daily Std Dev",    "1.15%"),
        ("Skewness",         "Negative (left tail)"),
        ("Excess Kurtosis",  "> 0 (leptokurtic)"),
        ("Is Leptokurtic?",  "Yes"),
    ]
    for ri, row in enumerate(stat_data):
        for ci, val in enumerate(row):
            cell = stat_tbl.cell(ri, ci)
            if ri == 0:
                set_cell_text(cell, val, bg_color=C_NAVY, font_color=C_WHITE,
                              font_size=9, bold=True)
            else:
                bg = C_LIGHT_GREY if ri % 2 == 1 else C_WHITE
                set_cell_text(cell, val, bg_color=bg, font_color=C_NAVY, font_size=9)

    add_textbox(sl, left_x + Inches(0.1), stat_top + Inches(1.55),
                col_w - Inches(0.2), Inches(0.55),
                "Fat tails: large losses occur more often than the normal model predicts.",
                9, italic=True, color=C_NAVY)

    # Per-ticker VaR table
    pticker_top = stat_top + Inches(2.2)
    add_textbox(sl, left_x + Inches(0.1), pticker_top - Inches(0.25),
                col_w - Inches(0.2), Inches(0.22),
                "Per-Ticker VaR Bounds", 9, bold=True, color=C_NAVY)
    pt_tbl = sl.shapes.add_table(5, 3,
                                  left_x + Inches(0.1),
                                  pticker_top,
                                  col_w - Inches(0.2),
                                  Inches(1.2)).table
    pt_data = [
        ("Ticker", "Dollar-Delta", "10-day VaR($)"),
        ("CBA",  "66,443",  "5,089"),
        ("BHP",  "58,338",  "4,770"),
        ("CSL",  "102,757", "10,407"),
        ("",     "Indep: $12,528 | Perf.corr: $20,265 | Empirical: $13,583", ""),
    ]
    for ri, row in enumerate(pt_data):
        for ci, val in enumerate(row):
            cell = pt_tbl.cell(ri, ci)
            if ri == 0:
                set_cell_text(cell, val, bg_color=C_NAVY, font_color=C_WHITE,
                              font_size=8, bold=True)
            elif ri == 4:
                tbl_merged = pt_tbl.cell(ri, 0)
                tbl_merged.merge(pt_tbl.cell(ri, 2))
                set_cell_text(tbl_merged, pt_data[4][1],
                              bg_color=RGBColor(0xCC, 0xD8, 0xFF),
                              font_color=C_NAVY, font_size=7)
                break
            else:
                bg = C_LIGHT_GREY if ri % 2 == 1 else C_WHITE
                set_cell_text(cell, val, bg_color=bg, font_color=C_NAVY, font_size=8)

    # RIGHT column
    add_rect(sl, right_x, content_y, col_w, content_h, RGBColor(0xE0, 0xE8, 0xFF))

    add_textbox(sl, right_x + Inches(0.1), content_y + Inches(0.08),
                col_w - Inches(0.2), Inches(0.35),
                "VaR Comparison", 12, bold=True, color=C_NAVY)

    right_items = [
        "Historical VaR ($11,825)  <  Parametric VaR ($13,583)",
        "This is unusual — normally hist > param for leptokurtic distributions.",
        "Our portfolio’s diversification across 3 sectors dampens tail extremity.",
        "Expected Shortfall ($21,431) is 81% above Historical VaR — confirming fat tails exist beyond the threshold.",
    ]
    r_y = content_y + Inches(0.45)
    for item in right_items:
        add_textbox(sl, right_x + Inches(0.1), r_y,
                    col_w - Inches(0.2), Inches(0.6),
                    "•  " + item, 10, color=C_NAVY)
        r_y += Inches(0.65)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Stress Correlation Analysis
# ─────────────────────────────────────────────────────────────────────────────
def slide11(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Stress Correlation Analysis",
                   "Multi-Asset MC VaR Across Three Correlation Regimes")
    add_slide_number(sl, 11)

    # Intro text
    add_textbox(sl, Inches(0.25), Inches(1.08), Inches(9.5), Inches(0.38),
                "Cross-sector correlations rise sharply in crises (GFC 2008, quant crash 2007), "
                "undermining diversification precisely when it is most needed.",
                9.5, color=C_NAVY)

    # Correlation matrix (compact, top-left)
    add_textbox(sl, Inches(0.25), Inches(1.5), Inches(3.5), Inches(0.25),
                "Empirical Correlation Matrix", 9, bold=True, color=C_NAVY)

    corr_tbl = sl.shapes.add_table(4, 4,
                                    Inches(0.25), Inches(1.77),
                                    Inches(3.2), Inches(0.8)).table
    corr_data = [
        ("",    "CBA", "BHP", "CSL"),
        ("CBA", "1.00", "0.07", "0.14"),
        ("BHP", "0.07", "1.00", "0.10"),
        ("CSL", "0.14", "0.10", "1.00"),
    ]
    for ri, row in enumerate(corr_data):
        for ci, val in enumerate(row):
            cell = corr_tbl.cell(ri, ci)
            if ri == 0 or ci == 0:
                set_cell_text(cell, val, bg_color=C_NAVY, font_color=C_WHITE,
                              font_size=8, bold=True)
            else:
                diag = C_ACCENT if ri == ci else None
                bg = diag if diag else (C_LIGHT_GREY if ri % 2 == 1 else C_WHITE)
                fc = C_WHITE if diag else C_NAVY
                set_cell_text(cell, val, bg_color=bg, font_color=fc, font_size=8)

    # Correlation stress chart (main body)
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart5_correlation_stress.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(0.15), Inches(2.65),
                                  Inches(9.7), Inches(2.75))

    # Key observation
    obs_y = Inches(5.15)
    add_textbox(sl, Inches(0.25), obs_y, Inches(9.5), Inches(0.28),
                "Key observation: A risk dashboard ignoring correlation regimes "
                "systematically underestimates crisis-period losses.",
                9, bold=True, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Scenario Analysis: Spot & Vol Shocks
# ─────────────────────────────────────────────────────────────────────────────
def slide12(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Scenario Analysis",
                   "Portfolio P&L Under Spot & Volatility Shocks (full repricing)")
    add_slide_number(sl, 12)

    # Compact table on top-right
    headers = ["Scenario", "P&L ($)", "P&L (%)"]
    data = [
        ("Base case",              "$0",        "0.00%"),
        ("Equity −10%",            "−$21,782",  "−7.11%"),
        ("Equity −20%",            "−$41,180",  "−13.43%"),
        ("Equity +10%",            "+$23,402",  "+7.64%"),
        ("Equity +20%",            "+$47,565",  "+15.52%"),
        ("Vol shock +30%",         "+$2,150",   "+0.70%"),
        ("Vol shock +50%",         "+$3,587",   "+1.17%"),
        ("Crash (−15%+50%vol)",    "≈−$22,000", "≈−7.2%"),
    ]

    tbl_left = Inches(6.35)
    tbl_top  = Inches(1.08)
    tbl_w    = Inches(3.5)
    tbl_h    = Inches(2.5)

    tbl = sl.shapes.add_table(len(data) + 1, 3,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(1.75), Inches(0.95), Inches(0.8)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=8, bold=True)

    for ri, row in enumerate(data):
        is_base  = ri == 0
        is_crash = ri == len(data) - 1
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if is_base:
                bg = C_NAVY; fc = C_WHITE; bd = True
            elif is_crash:
                bg = RGBColor(0xFF, 0xCC, 0xCC); fc = C_NAVY; bd = False
            else:
                bg = C_LIGHT_GREY if ri % 2 == 1 else C_WHITE
                fc = C_NAVY; bd = False
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=fc,
                          font_size=8, bold=bd, align=align)

    # Scenarios chart on the left
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart6_scenarios.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(0.15), Inches(1.08),
                                  Inches(6.05), Inches(3.7))

    # Observations below chart
    obs_y = Inches(4.85)
    add_textbox(sl, Inches(0.25), obs_y, Inches(9.5), Inches(0.5),
                "•  Equity shocks dominate P&L: ±20% spot ≈ ±$40–48k P&L\n"
                "•  Vol shocks alone are minimal — long CSL straddle offset by short BHP call",
                8.5, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Scenario Analysis: Yield Curve Shocks
# ─────────────────────────────────────────────────────────────────────────────
def slide13(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Scenario Analysis",
                   "Yield Curve Shocks — Rate Sensitivity Analysis")
    add_slide_number(sl, 13)

    headers = ["Scenario", "Short Shift", "Long Shift", "P&L ($)", "P&L (%)"]
    data = [
        ("Base case",                              "+0bp",   "+0bp",   "$0",      "0.000%"),
        ("Parallel rates +50bp",                   "+50bp",  "+50bp",  "−$197", "−0.064%"),
        ("Parallel rates +100bp",                  "+100bp", "+100bp", "−$392", "−0.128%"),
        ("Parallel rates −50bp",              "−50bp", "−50bp", "+$198", "+0.065%"),
        ("Steepener (−50bp short, +50bp long)", "−50bp", "+50bp", "small", "small"),
        ("Flattener (+50bp short, −50bp long)", "+50bp",  "−50bp", "small", "small"),
        ("Bear flattener (+100bp short, +25bp)",   "+100bp", "+25bp",  "≈−$295", "≈−0.10%"),
    ]

    tbl_left = Inches(0.15)
    tbl_top  = Inches(1.1)
    tbl_w    = Inches(9.7)
    tbl_h    = Inches(2.5)

    tbl = sl.shapes.add_table(len(data) + 1, 5,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(3.5), Inches(1.4), Inches(1.4), Inches(1.7), Inches(1.7)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=9, bold=True)

    for ri, row in enumerate(data):
        is_base = ri == 0
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if is_base:
                bg = C_NAVY; fc = C_WHITE; bd = True
            else:
                bg = C_LIGHT_GREY if ri % 2 == 1 else C_WHITE
                fc = C_NAVY; bd = False
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=fc,
                          font_size=9, bold=bd, align=align)

    # Observations
    obs_y = Inches(3.72)
    add_rect(sl, Inches(0.15), obs_y, Inches(9.7), Inches(1.6),
             RGBColor(0xE8, 0xED, 0xFF))
    add_textbox(sl, Inches(0.25), obs_y + Inches(0.05),
                Inches(9.4), Inches(0.28),
                "Key Observations", 10, bold=True, color=C_NAVY)
    obs_items = [
        "Rate sensitivity is modest (~$197–392 per 50–100bp) — no fixed-income positions, short-dated options",
        "Equity risk dominates; rate risk is second-order for this portfolio",
        "Rho = −$395 per 1% rate move; yield curve scenarios are consistent with per-position rho",
        "A real fixed-income overlay or duration-bearing position would change this significantly",
    ]
    obs_text = '\n'.join(f"•  {o}" for o in obs_items)
    add_textbox(sl, Inches(0.25), obs_y + Inches(0.36),
                Inches(9.4), Inches(1.1),
                obs_text, 9.5, color=C_NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Maximum Drawdown
# ─────────────────────────────────────────────────────────────────────────────
def slide14(prs, chart_dir=None):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_CONTENT)
    add_header_bar(sl, "Maximum Drawdown",
                   "Path-Dependent Portfolio Revaluation vs Equity Equivalent")
    add_slide_number(sl, 14)

    # MDD comparison table (compact, top)
    headers = ["Method", "MDD (%)", "Dollar MDD ($)", "Options?"]
    data = [
        ("Equity-equivalent (delta-weighted)",        "30.86%", "70,224", "No"),
        ("Full-portfolio (path-dependent repricing)", "11.89%", "36,460", "Yes"),
        ("Reduction",                                 "61.5%",  "33,764", "—"),
    ]

    tbl_left = Inches(0.25)
    tbl_top  = Inches(1.08)
    tbl_w    = Inches(9.5)
    tbl_h    = Inches(1.1)

    tbl = sl.shapes.add_table(len(data) + 1, 4,
                               tbl_left, tbl_top, tbl_w, tbl_h).table

    col_widths = [Inches(3.8), Inches(1.4), Inches(2.0), Inches(1.3)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_text(cell, h, bg_color=C_NAVY, font_color=C_WHITE,
                      font_size=10, bold=True)

    for ri, row in enumerate(data):
        is_red = ri == 2
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if is_red:
                bg = RGBColor(0x2D, 0x3A, 0x6B); fc = C_ACCENT; bd = True
            else:
                bg = C_LIGHT_GREY if ri % 2 == 0 else C_WHITE
                fc = C_NAVY; bd = False
            align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            set_cell_text(cell, val, bg_color=bg, font_color=fc,
                          font_size=10, bold=bd, align=align)

    # Drawdown chart (main body)
    if chart_dir:
        chart_path = os.path.join(chart_dir, 'chart7_drawdown.png')
        if os.path.exists(chart_path):
            sl.shapes.add_picture(chart_path,
                                  Inches(0.15), Inches(2.25),
                                  Inches(9.7), Inches(3.15))
    else:
        # Key insight fallback
        ki_y = Inches(2.3)
        add_rect(sl, Inches(0.25), ki_y, Inches(9.5), Inches(1.0),
                 RGBColor(0x2D, 0x3A, 0x6B))
        add_textbox(sl, Inches(0.35), ki_y + Inches(0.08),
                    Inches(1.5), Inches(0.3),
                    "Key Insight:", 11, bold=True, color=C_ACCENT)
        add_textbox(sl, Inches(0.35), ki_y + Inches(0.38),
                    Inches(9.1), Inches(0.55),
                    "The option overlays reduce historical drawdown by 61.5% — observed in the historical price path.",
                    10, color=C_WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — Limitations & Extensions
# ─────────────────────────────────────────────────────────────────────────────
def slide15(prs):
    sl = blank_slide(prs)
    set_bg(sl, C_BG_LAST)
    add_slide_number(sl, 15)

    # Title
    add_textbox(sl, Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.65),
                "Limitations & Extensions",
                28, bold=True, color=C_WHITE)

    # Subtle accent bar
    add_rect(sl, Inches(0), Inches(0.95), Inches(10), Inches(0.04), C_ACCENT)

    col_w = Inches(4.55)
    col_y = Inches(1.05)
    col_h = Inches(4.35)

    # LEFT column — LIMITATIONS
    add_rect(sl, Inches(0.2), col_y, col_w, col_h, RGBColor(0x26, 0x33, 0x7A))

    add_textbox(sl, Inches(0.35), col_y + Inches(0.1), col_w - Inches(0.3), Inches(0.35),
                "LIMITATIONS", 14, bold=True, color=C_ACCENT)

    limitations = [
        ("1.  Sqrt-time VaR scaling",
         "Assumes i.i.d. returns. Real returns exhibit volatility clustering. "
         "Effect: underestimates multi-day VaR in stressed regimes."),
        ("2.  Static volatility (no smile)",
         "Black-Scholes assumes constant σ. Real vol is time-varying and mean-reverting. "
         "2-year historical σ may lag regime changes."),
        ("3.  Single-underlying MC VaR",
         "Portfolio contains three underlyings; the single-asset MC in "
         "Portfolio.monte_carlo_var() misses cross-asset correlation. "
         "Addressed in Section 7 via multi-asset MC."),
    ]

    lim_y = col_y + Inches(0.5)
    for label, body in limitations:
        tb_l = sl.shapes.add_textbox(Inches(0.35), lim_y, col_w - Inches(0.3), Inches(0.28))
        tf = tb_l.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = 'Calibri'
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = C_ACCENT

        add_textbox(sl, Inches(0.35), lim_y + Inches(0.26), col_w - Inches(0.3), Inches(0.85),
                    body, 9, color=RGBColor(0xCC, 0xD6, 0xFF))
        lim_y += Inches(1.18)

    # RIGHT column — FUTURE EXTENSIONS
    right_x = Inches(5.05)
    add_rect(sl, right_x, col_y, col_w, col_h, RGBColor(0x26, 0x33, 0x7A))

    add_textbox(sl, right_x + Inches(0.15), col_y + Inches(0.1), col_w - Inches(0.3), Inches(0.35),
                "FUTURE EXTENSIONS", 14, bold=True, color=C_ACCENT)

    extensions = [
        "GARCH(1,1) model for time-varying volatility and better multi-day VaR scaling",
        "Volatility smile / local vol surface (Dupire) for OTM option pricing accuracy",
        "Fixed-income overlay (interest rate swaps, bond futures) to manage rate risk",
        "Real-time data pipeline — replace cached CSV with live RBA / Bloomberg feed",
        "Stress-testing framework aligned with APRA/Basel III 10-day 99% VaR standard",
    ]

    ext_y = col_y + Inches(0.52)
    for i, ext in enumerate(extensions):
        # Number circle simulation
        add_textbox(sl, right_x + Inches(0.1), ext_y, Inches(0.35), Inches(0.38),
                    str(i + 1) + ".", 11, bold=True, color=C_ACCENT)
        add_textbox(sl, right_x + Inches(0.45), ext_y, col_w - Inches(0.6), Inches(0.65),
                    ext, 9.5, color=RGBColor(0xCC, 0xD6, 0xFF))
        ext_y += Inches(0.73)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def main():
    chart_dir = '/tmp/pptx_charts'
    print("Generating charts …")
    generate_charts(chart_dir)

    prs = new_prs()

    print("Building slide 1  — Title Slide …")
    slide1(prs)
    print("Building slide 2  — Platform Overview …")
    slide2(prs)
    print("Building slide 3  — Yield Curve Construction …")
    slide3(prs, chart_dir)
    print("Building slide 4  — Interpolation Comparison …")
    slide4(prs)
    print("Building slide 5  — Derivatives Pricing …")
    slide5(prs, chart_dir)
    print("Building slide 6  — Option Greeks …")
    slide6(prs)
    print("Building slide 7  — Portfolio Construction …")
    slide7(prs)
    print("Building slide 8  — Portfolio Exposures …")
    slide8(prs, chart_dir)
    print("Building slide 9  — Value at Risk …")
    slide9(prs, chart_dir)
    print("Building slide 10 — Leptokurtosis …")
    slide10(prs)
    print("Building slide 11 — Stress Correlation …")
    slide11(prs, chart_dir)
    print("Building slide 12 — Scenario: Spot & Vol …")
    slide12(prs, chart_dir)
    print("Building slide 13 — Scenario: Yield Curve …")
    slide13(prs)
    print("Building slide 14 — Maximum Drawdown …")
    slide14(prs, chart_dir)
    print("Building slide 15 — Limitations & Extensions …")
    slide15(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\nSUCCESS: Saved to {OUT_PATH}")
    print(f"         Slides : {len(prs.slides)}")
    print(f"         Size   : {os.path.getsize(OUT_PATH):,} bytes")


if __name__ == "__main__":
    main()
